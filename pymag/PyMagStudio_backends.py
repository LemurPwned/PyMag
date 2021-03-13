#todo dodać backend ze zdalnego pulpitu
import datetime
from pymag.engine.solver import SimulationRunner
import queue as Queue
from threading import Thread

import cmtj
from PyQt5.QtWidgets import QFileDialog, QWidget

import os.path
import pickle
import sys
import time

import numpy as np
from natsort import natsorted
from pyqtgraph.dockarea import Dock, DockArea

from pymag.engine import task_runner
from pymag.gui.core import *

stop = 1


class measurement_Menegement():
    def __init__(self, parent):
        super(measurement_Menegement, self).__init__()
        self.active_experiment = None
        self.active_experiments = []
        self.table_experiments = pg.TableWidget(editable=False, sortable=False)
        self.simulations_list = {
            "results": [],
            "settings": [],
            "layer_params": [],
            "simulation_params": []
        }
        self.RemoveButton = QtWidgets.QPushButton()
        self.RemoveButton.setText("Remove selected result")
        self.RemoveButton.clicked.connect(self.removeLayer)
        self.RemoveButton.setEnabled(False)
        self.ExportSelected = QtWidgets.QPushButton()
        self.ExportSelected.setText("Export selected to .csv")
        self.ExportSelected.clicked.connect(self.exportSelected)
        self.ExportSelected.setEnabled(False)
        self.ctrlWidget = QtGui.QWidget()
        self.ctrLayout = QtGui.QGridLayout()
        self.ctrlWidget.setLayout(self.ctrLayout)
        self.ctrLayout.addWidget(self.table_experiments)
        self.ctrLayout.addWidget(self.RemoveButton)
        self.ctrLayout.addWidget(self.ExportSelected)
        self.table_experiments.cellDoubleClicked.connect(self.clicked2x)

    def removeLayer(self):
        for n in self.active_experiments:
            self.simulations_list["settings"].pop(n)
            self.simulations_list["results"].pop(n)
            self.simulations_list["layer_params"].pop(n)
            self.simulations_list["simulation_params"].pop(n)
            self.table_experiments.setData(self.simulations_list["settings"])
        self.active_experiments = []
        plotter.replot_results()
        self.print_and_color_table()

    def exportSelected(self):
        plotter.replot_results(self.active_experiment, save=1)

    def clicked2x(self, parent):
        n = self.table_experiments.currentRow()
        m = int(self.table_experiments.rowCount())
        if n in self.active_experiments:
            self.active_experiments.remove(n)
        else:
            self.active_experiments.append(n)
        if not self.active_experiments:
            self.RemoveButton.setEnabled(False)
            self.ExportSelected.setEnabled(False)
        else:
            self.RemoveButton.setEnabled(True)
            self.ExportSelected.setEnabled(True)

        for i in range(0, m):
            if i in self.active_experiments:
                self.simulations_list["settings"][i][0] = "V"
            else:
                self.simulations_list["settings"][i][0] = "X"
        self.print_and_color_table()
        plotter.replot_results()

    def print_and_color_table(self):
        m = int(self.table_experiments.rowCount())
        self.table_experiments.setData(self.simulations_list["settings"])
        self.table_experiments.setHorizontalHeaderLabels(
            ["Select", "Type", "Timestamp"])
        for i in range(0, m):
            if i in self.active_experiments:
                self.table_experiments.item(i, 0).setBackground(
                    QtGui.QColor(255, 0, 0))  ##dosnt work!
            else:
                self.table_experiments.item(i, 0).setBackground(
                    QtGui.QColor(255, 255, 255))


class Ui_MainWindow(
        QWidget, ):
    def __init__(self):
        super().__init__()

        #defaults file names and load
        self.defaultStimulusFile = "defaultStimulus.csv"
        self.defaultParametersFile = "defaultParameters.csv"
        self.loadDefaults()

        MainWindow.setObjectName(PyMagVersion)
        MainWindow.setWindowTitle(PyMagVersion)
        MainWindow.resize(1200, 900)

        #dock area as a central widget of GUI
        self.area = DockArea()
        MainWindow.setCentralWidget(self.area)

        # define classes
        self.WindowSettings = Settings(self)
        self.WindowAbout = About(self)

        self.LineShape = Line_shape()
        self.SD_plot = plotDynamics()
        self.PIMM_plot = plotDynamics()
        self.ResPlot = Res_plot()
        self.MagPlot = Mag_plot()
        self.ctrLayout = addMenuBar(self)
        self.ctrWidget = self.ctrLayout.ctrlWidget
        self.Trajectory_plot = TrajectoryPlot()
        self.simulationsMenegement = measurement_Menegement(self)
        self.measurementsMenegement = measurement_Menegement(self)
        self.widget_layer_params = paramsAndStimulus(self)
        self.table_results = pg.TableWidget(editable=True, sortable=False)
        self.table_results.setHorizontalHeaderLabels(ResultsColumns)

        #define dock area structure
        """
        Maybe loop instead of declation

        for all of the below

        d1 - d10 

        """
        self.d1 = Dock("Control panel", size=(300, 50))
        self.d2 = Dock("PIMM-FMR", size=(300, 300))
        self.d7 = Dock("SD-FMR", size=(300, 300))
        self.d3 = Dock("Magnetization", size=(300, 300))
        self.d4 = Dock("Simulation results", size=(300, 200))
        self.d6 = Dock("Resistance", size=(300, 300))
        self.d10 = Dock("Trajectory", size=(200, 200))
        self.d11 = Dock("SD lines", size=(300, 300))
        self.d20 = Dock("Layer parameters", size=(100, 200))
        self.d13 = Dock("Python console", size=(200, 200))
        self.d14 = Dock("Simulation management", size=(200, 200))
        self.d15 = Dock("Measurement management", size=(200, 200))

        self.area.addDock(self.d1, 'left')
        self.area.addDock(self.d2, 'right')
        self.area.addDock(self.d3, 'bottom', self.d1)
        self.area.moveDock(self.d15, 'right', self.d2)
        self.area.moveDock(self.d14, 'above', self.d15)
        self.area.addDock(self.d4, 'above', self.d3)
        self.area.addDock(self.d6, 'above', self.d3)
        self.area.addDock(self.d7, 'above', self.d2)
        self.area.addDock(self.d10, 'above', self.d7)
        self.area.addDock(self.d11, 'above', self.d7)
        self.area.moveDock(self.d20, 'top', self.d2)

        self.d20.addWidget(self.widget_layer_params.ctrlWidget)
        self.d14.addWidget(self.simulationsMenegement.ctrlWidget)
        self.d15.addWidget(self.measurementsMenegement.ctrlWidget)
        self.d11.addWidget(self.LineShape.plotsLS)
        self.d7.addWidget(self.SD_plot.plotsDynamics_view)
        self.d2.addWidget(self.PIMM_plot.plotsDynamics_view)
        self.d6.addWidget(self.ResPlot.plotsRes)
        self.d3.addWidget(self.MagPlot.plotsMag)
        self.d4.addWidget(self.table_results)
        self.d1.addWidget(self.ctrWidget)
        self.d10.addWidget(self.Trajectory_plot.w)

        self.ports = []

        self.timer = pg.QtCore.QTimer()
        self.timer.timeout.connect(self.update)
        self.timer.start(0)

        MainWindow.show()

        self.p = Thread(target=Simulate)
        self.p.start()

    def loadDefaults(self):
        try:
            self.StimulusParameters = pd.read_csv(self.defaultStimulusFile,
                                                  delimiter='\t')
        except:
            print("No file with default stimulus was found")
            self.StimulusParameters = pd.DataFrame(
                np.array([[
                    -800000, 800000, 50, 0, 89.9, 45, 0, 47e9, 48, 0.01, 0.01,
                    "[1 0 0]", 0, "[1 0 0]", 4e-9, 2000
                ]]),
                columns=[
                    'Hmin', 'Hmax', 'Hsteps', 'Hback', 'HTheta', 'HPhi',
                    'fmin', 'fmax', 'fsteps', 'IAC', 'IDC', 'Idir', 'fphase',
                    'Vdir', 'LLGtime', 'LLGsteps'
                ])

        try:
            self.layerParameters = pd.read_csv(self.defaultParametersFile,
                                               delimiter='\t')
        except:
            print("No file with default parameters was found")
            self.layerParameters = pd.DataFrame(
                np.array([[
                    1, 1.6, 3000, "[1 0 0]", -1e-5, 0.01, 1e-9, "[0 0 1]",
                    0.02, 0.01, 0.01, 100, 120
                ],
                          [
                              2, 1.1, 4000, "[1 0 0]", -1e-5, 0.01, 1e-9,
                              "[0 0 1]", 0.02, 0.01, 0.01, 100, 120
                          ]]),
                columns=[
                    'layer', 'Ms', 'Ku', 'Kdir', 'J', 'alpha', 'th', 'N',
                    'AMR', 'SMR', 'AHE', 'Rx0', 'Ry0'
                ])

    def newSettings(self):
        self.WindowSettings.signal.connect(self.slot)
        self.WindowSettings.show()

    def newAbout(self):
        self.WindowAbout.show()

    def fullScreenMode(self):
        if MainWindow.isFullScreen():
            MainWindow.showNormal()
        else:
            MainWindow.showFullScreen()

    def endProgram(self):
        self.getDfFromTable(
            self.widget_layer_params.table_stimulus_params).to_csv(
                self.defaultStimulusFile,
                encoding='utf-8',
                index=False,
                sep='\t')
        self.getDfFromTable(
            self.widget_layer_params.table_layer_params).to_csv(
                self.defaultParametersFile,
                encoding='utf-8',
                index=False,
                sep='\t')
        os._exit(0)

    def setStimulusForALl(self):
        df_stimulus = self.getDfFromTable(
            self.widget_layer_params.table_stimulus_params)
        for n in range(
                self.simulationsMenegement.table_experiments.rowCount()):
            plotter.simulationsMenegement.simulations_list[
                "simulation_params"][n] = df_stimulus
        plotter.replot_results()

    def clearPlot(self):
        self.ResPlot.clearPlot()
        self.MagPlot.clearPlot()
        self.PIMM_plot.clearPlot()
        self.SD_plot.clearPlot()
        self.LineShape.clearPlot()
        self.Trajectory_plot.clear()

    def OpenDailogBox(self, extention=".csv"):
        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        fileName, _ = QFileDialog.getOpenFileName(
            self,
            "QFileDialog.getOpenFileName()",
            "",
            "File (*" + extention + ");;CSV (*.csv)",
            options=options)
        return fileName

    def OpenDirecoryDailogBox(self):
        dialog = QtWidgets.QFileDialog()
        dialog.setFileMode(QtWidgets.QFileDialog.Directory)
        dialog.setOption(QtWidgets.QFileDialog.DontUseNativeDialog, True)
        fileName, _ = QFileDialog.getOpenFileName(
            self, "QFileDialog.getOpenFileName()", "")
        return fileName

    def SaveDailogBox(self, extention=".csv"):
        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        fileName, _ = QFileDialog.getSaveFileName(
            self,
            "QFileDialog.getOpenFileName()",
            "",
            "File (*" + extention + ")",
            options=options)
        return fileName + extention

    def saveParams(self, auto=0):
        if auto == 0:
            fileName = self.SaveDailogBox()
        else:
            curr_dir = os.path.dirname(os.path.realpath(__file__))
            fileName = curr_dir + "/" + "previous_params.csv"
        if fileName:
            df_generated_data = self.getDfFromTable(
                self.widget_layer_params.table_layer_params)
            df_generated_data.to_csv(fileName,
                                     encoding='utf-8',
                                     index=False,
                                     sep='\t')

    def saveBinary(self):
        fileName = self.SaveDailogBox(".bin")
        a_file = open(fileName, "wb")
        pickle.dump([
            self.simulationsMenegement.simulations_list,
            self.measurementsMenegement.simulations_list
        ], a_file)
        a_file.close()

    def loadBinary(self):
        fileName = self.OpenDailogBox(".bin")
        a_file = open(fileName, "rb")
        package1 = pickle.load(a_file)
        self.simulationsMenegement.simulations_list = package1[0]
        self.measurementsMenegement.simulations_list = package1[1]
        a_file.close()
        plotter.simulationsMenegement.table_experiments.setData(
            pd.DataFrame(
                np.array(self.simulationsMenegement.
                         simulations_list["settings"])).to_numpy())
        plotter.measurementsMenegement.table_experiments.setData(
            pd.DataFrame(
                np.array(self.measurementsMenegement.
                         simulations_list["settings"])).to_numpy())

    def addToSimulationList(self):
        print("IM CALLED")
        df = self.getDfFromTable(self.widget_layer_params.table_layer_params)
        df_stimulus = self.getDfFromTable(
            self.widget_layer_params.table_stimulus_params)
        plotter.simulationsMenegement.simulations_list["results"].append({
            "MR":
            np.array([[0, 0, 0, 0, 0, 0], [1, 0, 0, 0, 0, 0]]),
            "SD_freqs":
            np.array([0, 1]),
            "SD":
            np.array([[0, 0, 0], [0, 0, 0]]),
            "PIMM_freqs":
            1,
            "PIMM":
            np.array([[0, 0]]),
            "traj":
            np.array([[0, 0, 0]])
        })
        plotter.simulationsMenegement.simulations_list["settings"].append(
            ["X", "from_table", "To be simulated"])
        plotter.simulationsMenegement.simulations_list["layer_params"].append(
            df)
        plotter.simulationsMenegement.simulations_list[
            "simulation_params"].append(df_stimulus)
        plotter.simulationsMenegement.print_and_color_table()
        plotter.simulationsMenegement.print_and_color_table()

    def loadMultipleLayerParams(self):
        dirName = "/home/sz/Desktop/AGH/PyMag/20 Oct 2020/4651_Pymag"
        listaa = []
        listab = []
        listac = []
        for path, subdirs, files in os.walk(dirName):

            for name in files:
                if os.path.splitext(name)[1] == ".csv":
                    listaa.append(os.path.join(path, name))
                elif os.path.splitext(name)[1] == ".dat" and os.path.splitext(
                        name)[0].find("PO") != -1:
                    listab.append(os.path.join(path, name))
                elif os.path.splitext(name)[1] == ".dat" and os.path.splitext(
                        name)[0].find("ment") != -1:
                    listac.append(os.path.join(path, name))
        for a in natsorted(listaa):
            self.loadParams(auto=False, filename=a)
        for b in natsorted(listab):
            self.loadResults(auto=False, filename=b)
        for c in natsorted(listac):
            self.loadResults(auto=False, filename=c)

    def loadParams(self, dialog=True, auto=True, filename=""):
        if auto == True:
            fileName = self.OpenDailogBox()
        else:
            fileName = filename
        if fileName:
            df = pd.read_csv(fileName, delimiter='\t')
            self.widget_layer_params.table_layer_params.setData(np.array(df))
            self.widget_layer_params.table_layer_params.setHorizontalHeaderLabels(
                df.columns)
            simulationName = os.path.basename(fileName).replace(
                ".csv",
                "").replace("_layer_params",
                            "").replace("_layers_params",
                                        "").replace("_layer_param", "")
            self.ctrLayout.Simulation_Name.setText(simulationName)
            plotter.simulationsMenegement.simulations_list["results"].append({
                "MR":
                np.array([[0, 0, 0, 0, 0, 0], [1, 0, 0, 0, 0, 0]]),
                "SD_freqs":
                np.array([0, 1]),
                "SD":
                np.array([[0, 0, 0], [0, 0, 0]]),
                "PIMM_freqs":
                1,
                "PIMM":
                np.array([[0, 0]]),
                "traj":
                np.array([[0, 0, 0]])
            })
            plotter.simulationsMenegement.simulations_list["settings"].append(
                ["X", simulationName, "To be simulated"])
            plotter.simulationsMenegement.simulations_list[
                "layer_params"].append(df)
            df_stimulus = self.getDfFromTable(
                self.widget_layer_params.table_stimulus_params)
            plotter.simulationsMenegement.simulations_list[
                "simulation_params"].append(df_stimulus)
            plotter.simulationsMenegement.print_and_color_table()
            plotter.simulationsMenegement.print_and_color_table()

    def getDfFromTable(self, table):
        number_of_rows = table.rowCount()
        number_of_columns = table.columnCount()
        tmp_df = pd.DataFrame()
        tmp_col_name = []
        for i in range(number_of_columns):
            tmp_col_name.append(table.takeHorizontalHeaderItem(i).text())
            for j in range(number_of_rows):
                tmp_df.loc[j, i] = table.item(j, i).text()
        table.setHorizontalHeaderLabels(tmp_col_name)
        tmp_df.columns = tmp_col_name
        return tmp_df

    def saveResults(self):
        fileName = self.SaveDailogBox("_MR.csv")
        if fileName:
            df_generated_data = self.getDfFromTable(self.table_results)
            df_generated_data.to_csv(fileName,
                                     encoding='utf-8',
                                     index=False,
                                     sep='\t')

    def appendPptx(self, mode=0, reportName="test.docx"):
        """
        To powinien byc oddzielny modul 
        """
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        slide_size = (24, 36)
        prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(
            slide_size[1])

        grid_size = 2

        i = 0

        def df_to_pptx_table(df, l, t, w, h):

            df = df.iloc[[0, 1], [1, 2, 3, 4, 5, 6, 7]]

            rows = df.shape[0] + 1
            cols = df.shape[1]

            left = l
            top = t
            width = w
            height = h

            shape = slide.shapes.add_table(rows, cols, left, top, width,
                                           height)
            table = shape.table

            tbl = shape._element.graphic.graphicData.tbl
            style_id = "{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}"
            tbl[0][-1].text = style_id

            for i in range(0, cols):
                table.cell(0, i).text = df.columns[i]
                for j in range(1, rows):
                    table.cell(j, i).text = str(df.iloc[j - 1, i])
            """
            To nie moze tak byc import w srodku!
            chyba ze dodany na stale w requirements
            """

            from pptx.util import Pt

            def iter_cells(table):
                for row in table.rows:
                    for cell in row.cells:
                        yield cell

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)

        for num in range(
                0,
                int(plotter.simulationsMenegement.table_experiments.rowCount())
        ):
            self.export_images(num)
            left = Inches(grid_size)
            top = Inches(grid_size + i * grid_size + 0.1)
            top_text = Inches(grid_size + grid_size / 2 + i * grid_size + 0.1)
            h = Inches(grid_size)
            w = Inches(grid_size)
            shape = slide.shapes.add_textbox(0, top_text, w, h)
            text_frame = shape.text_frame
            Simulation_Name = plotter.simulationsMenegement.simulations_list[
                "settings"][num][1]
            text_frame.text = Simulation_Name
            text_frame.word_wrap = False
            slide.shapes.add_picture('Mag.png', left, top, w, h)
            slide.shapes.add_picture('Res.png', 2 * left, top, w, h)
            slide.shapes.add_picture('SD.png', 3 * left, top, w, h)
            slide.shapes.add_picture('PIMM.png', 4 * left, top, w, h)
            slide.shapes.add_picture('LS.png', 5 * left, top, w, h)
            i = i + 1
            df_to_pptx_table(
                self.getDfFromTable(
                    self.widget_layer_params.table_layer_params), 6 * left,
                top, 5 * w, Inches(grid_size / 2))
        prs.save('test.pptx')

    def export_images(self, num=1):
        plotter.simulationsMenegement.active_experiments = [num]
        plotter.measurementsMenegement.active_experiments = [
            num, num +
            int(plotter.simulationsMenegement.table_experiments.rowCount())
        ]
        """
        Loop ponizej zamiast iteracji
        """
        plotter.replot_results()
        exporter = pg.exporters.ImageExporter(self.ResPlot.plotsRes.sceneObj)
        exporter.export('Res' + '.png')
        exporter = pg.exporters.ImageExporter(self.MagPlot.plotsMag.sceneObj)
        exporter.export('Mag' + '.png')
        exporter = pg.exporters.ImageExporter(
            self.SD_plot.plotsDynamics_view.sceneObj)
        exporter.export('SD' + '.png')
        # self.PIMM_plot.export()
        exporter = pg.exporters.ImageExporter(
            self.PIMM_plot.plotsDynamics_view.sceneObj)
        exporter.export('PIMM' + '.png')
        exporter = pg.exporters.ImageExporter(self.LineShape.plotsLS.sceneObj)
        exporter.export('LS' + '.png')

    def saveReport(self, mode=0, reportName="test.docx"):

        if mode == 0:
            fileName = self.SaveDailogBox(".docx")
            print(fileName)
        else:

            import os
            curr_dir = os.path.dirname(os.path.realpath(__file__))
            fileName = curr_dir + "/" + reportName
            # powinno być! Bo te ściezki nie zadziałają na Windows!
            # filename = os.path.join(curr_dir, report_name)
            print(fileName)

        if fileName:
            """
            Podwójne importy! tez powinno byc na wierzchu
            """

            from datetime import datetime

            from docx import Document
            from docx.shared import Pt
            now = datetime.now()
            import pyqtgraph.exporters  # to juz bylo gdzies indziej importowane
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            dt_string = now.strftime("%d.%m.%Y %H:%M:%S")

            self.export_images()

            if mode == 0:
                document = Document()
            else:
                from pathlib import Path

                if Path(fileName).is_file():
                    document = Document(fileName)
                else:
                    document = Document()

            document.add_heading("Sample name", 0)
            # header section
            header_section = document.sections[0]
            header = header_section.header
            header_text = header.paragraphs[0]
            # header_text.text = 'Report prepared in ' + PyMagVersion + ' in version from ' + PyMagDate
            header_text.text = "Report timestamp: " + dt_string
            header_text.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            header_text.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            document.add_paragraph('Sample layer structure',
                                   style='List Number')

            def add_doc_table(df):

                data1 = self.getDfFromTable(df)
                table = document.add_table(rows=1, cols=len(data1.columns))
                hdr_cells = table.rows[0].cells
                for i, cl in enumerate(hdr_cells):
                    hdr_cells[i].text = data1.columns[i]
                    # hdr_cells[i]._tc.tcPr.tcW.type = 'auto'
                for row in data1.index:
                    values = data1.loc[row]
                    row_cells = table.add_row().cells
                    for i, v in enumerate(values):
                        row_cells[i].text = v

                for row in table.rows:
                    for cell in row.cells:
                        paragraphs = cell.paragraphs
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                font = run.font
                                font.size = Pt(6)

            add_doc_table(self.widget_layer_params.table_layer_params)
            picture_width_mm = 140
            """
            Petla ponizej
            """

            document.add_paragraph('Simulation of resistance',
                                   style='List Number')
            document.add_picture('Mag.png', width=36000 * picture_width_mm)
            document.add_paragraph('Simulation of magnetization',
                                   style='List Number')
            document.add_picture('Res.png', width=36000 * picture_width_mm)
            document.add_paragraph('Magnetization dynamics',
                                   style='List Number')
            document.add_picture('SD.png', width=36000 * picture_width_mm)
            document.add_picture('PIMM.png', width=36000 * picture_width_mm)
            document.add_picture('LS.png', width=36000 * picture_width_mm)
            document.add_paragraph('Numerical results', style='List Number')
            document.save(fileName)

    def loadResults(self, dialog=True, auto=True, filename=""):
        if auto == True:
            fileName = self.OpenDailogBox(".dat")
        else:
            fileName = filename
        if fileName:
            df = pd.read_csv(fileName, delimiter='\t')
            self.plotExperimentalData(
                df,
                os.path.basename(fileName).replace(".dat", ""))

    def plotExperimentalData(self, df, fileName):
        """
        Petla poniej i tylko jeden try catch

        Nie moze byc try catch: pass! 
        musi łapać konkrenty błąd!

        """

        try:
            self.MagPlot.Mx.plot(df['H'],
                                 df['Mx'],
                                 symbol='o',
                                 pen=None,
                                 symbolPen=(255, 0, 0, 90),
                                 symbolBrush=(255, 0, 0, 50))
        except:
            pass
        try:
            self.MagPlot.My.plot(df['H'],
                                 df['My'],
                                 symbol='o',
                                 pen=None,
                                 symbolBrush=(0, 255, 0, 90),
                                 symbolPen=(0, 255, 0, 50))
        except:
            pass
        try:
            self.MagPlot.Mz.plot(df['H'],
                                 df['Mz'],
                                 symbol='o',
                                 pen=None,
                                 symbolBrush=(0, 0, 255, 90),
                                 symbolPen=(0, 0, 255, 50))
        except:
            pass
        try:
            self.ResPlot.Rx.plot(df['H'],
                                 df['Rx'],
                                 symbol='o',
                                 pen=None,
                                 symbolBrush=(255, 0, 0, 90),
                                 symbolPen=(255, 0, 0, 50))
        except:
            pass
        try:
            self.ResPlot.Ry.plot(df['H'],
                                 df['Ry'],
                                 symbol='o',
                                 pen=None,
                                 symbolBrush=(0, 255, 0, 90),
                                 symbolPen=(0, 255, 255, 50))
        except:
            pass
        try:
            self.ResPlot.Rz.plot(df['H'],
                                 df['Rz'],
                                 symbol='o',
                                 pen=None,
                                 symbolBrush=(0, 0, 255, 90),
                                 symbolPen=(0, 0, 255, 50))

            plotter.measurementsMenegement.simulations_list["settings"].append(
                ["X", "Exp", fileName])
            plotter.measurementsMenegement.simulations_list["results"].append({
                "H":
                df['H'],
                "f":
                df['f']
            })
            plotter.measurementsMenegement.simulations_list[
                "layer_params"].append(["X"])
            plotter.measurementsMenegement.print_and_color_table()
        except:
            pass

    def btn_clk(self):
        self.ctrLayout.progress.setValue(0)
        global stop
        stop = 0

    def stop_clk(self):
        global stop
        print(stop)
        stop = 1

    def saveDockState(self):
        global state
        state = self.area.saveState()

    def loadDockState(self):
        global state
        self.area.restoreState(state)

    def getPort(self):
        q = Queue.Queue()
        self.ports.append((q))
        return q

    def replot_experimental(self):
        try:
            for k in self.measurementsMenegement.active_experiments:
                xp1 = self.measurementsMenegement.simulations_list["results"][
                    k]
                if "PO" in str(self.measurementsMenegement.
                               simulations_list["settings"][k][2]):  ###Raport
                    self.PIMM_plot.plotsSpectrum.plot(xp1['H'],
                                                      xp1['f'],
                                                      symbol='o',
                                                      pen=None,
                                                      symbolBrush=(255, 0, 0,
                                                                   70),
                                                      symbolPen=(255, 0, 0,
                                                                 70))
                    self.SD_plot.plotsSpectrum.plot(xp1['H'],
                                                    xp1['f'],
                                                    symbol='o',
                                                    pen=None,
                                                    symbolBrush=(255, 0, 0,
                                                                 70),
                                                    symbolPen=(255, 0, 0, 70))
                elif "ment" in str(
                        self.measurementsMenegement.
                        simulations_list["settings"][k][2]):  ###Raport
                    self.PIMM_plot.plotsSpectrum.plot(xp1['H'],
                                                      xp1['f'],
                                                      symbol='o',
                                                      pen=None,
                                                      symbolBrush=(0, 255, 0,
                                                                   120),
                                                      symbolPen=(0, 255, 0,
                                                                 120))
                    self.SD_plot.plotsSpectrum.plot(xp1['H'],
                                                    xp1['f'],
                                                    symbol='o',
                                                    pen=None,
                                                    symbolBrush=(0, 255, 0,
                                                                 120),
                                                    symbolPen=(0, 255, 0, 120))
                else:
                    self.PIMM_plot.plotsSpectrum.plot(xp1['H'],
                                                      xp1['f'],
                                                      symbol='o',
                                                      pen=None,
                                                      symbolBrush=(255, 0, 0,
                                                                   70),
                                                      symbolPen=(255, 0, 0,
                                                                 70))
                    self.SD_plot.plotsSpectrum.plot(xp1['H'],
                                                    xp1['f'],
                                                    symbol='o',
                                                    pen=None,
                                                    symbolBrush=(255, 0, 0,
                                                                 70),
                                                    symbolPen=(255, 0, 0, 70))
        except:
            pass

    def replot_all(self, xp, plot_realtime=0, save=0):
        try:
            """
            Tutaj poniej trzeba jakoś okomentować 
            """
            mode = xp["mode"]
            H = xp["MR"][:, 0]
            Mx = xp["MR"][:, 1]
            My = xp["MR"][:, 2]
            Mz = xp["MR"][:, 3]
            Rx = xp["MR"][:, 4]
            Ry = xp["MR"][:, 5]
            Rz = xp["MR"][:, 6]
            SD_freqs = xp["SD_freqs"]
            PIMM_delta_f = xp["PIMM_freqs"]
            if plot_realtime == 0:
                self.PIMM_plot.init_setup(Xmin=min(H),
                                          Xmax=max(H),
                                          Xsteps=len(H),
                                          dy=PIMM_delta_f)
                self.SD_plot.init_setup(Xmin=min(H),
                                        Xmax=max(H),
                                        Xsteps=len(H),
                                        dy=SD_freqs[1] - SD_freqs[0])
            self.MagPlot.Mx.plot(H, Mx, pen=(255, 0, 0))
            self.MagPlot.My.plot(H, My, pen=(0, 255, 0))
            self.MagPlot.Mz.plot(H, Mz, pen=(0, 0, 255))
            self.MagPlot.Mx.setYRange(-1.5, 1.5, padding=0)
            self.MagPlot.My.setYRange(-1.5, 1.5, padding=0)
            self.MagPlot.Mz.setYRange(-1.5, 1.5, padding=0)
            self.MagPlot.setMode(mode)
            self.ResPlot.Rx.plot(H, Rx, pen=(255, 0, 0))
            self.ResPlot.Ry.plot(H, Ry, pen=(0, 255, 0))
            self.ResPlot.Rz.plot(H, Rz, pen=(0, 0, 255))
            self.ResPlot.setMode(mode)

            self.SD_plot.update(xp["SD"],
                                SD_freqs[1] - SD_freqs[0],
                                H,
                                rm_bkg=1)
            self.SD_plot.setMode(mode)

            self.PIMM_plot.update(xp["PIMM"], PIMM_delta_f, H)
            self.PIMM_plot.setMode(mode)

            self.LineShape.update(xp["SD"], SD_freqs[1] - SD_freqs[0], H)
            self.LineShape.setMode(mode)

            self.Trajectory_plot.pltTraj(
                xp["traj"],
                [10 * len(H), 255 - 10 * len(H), 255 - 10 * len(H)])
            self.table_results.setData(xp["MR"].astype(np.str))
            if save == 1:
                fileName = self.SaveDailogBox("")
                if fileName:
                    pd.DataFrame(
                        xp["SD"].T,
                        columns=np.array(H)).set_index(SD_freqs).to_csv(
                            fileName + "_SD.csv",
                            encoding='utf-8',
                            index=True,
                            sep='\t')
        except:
            pass

    def replot_results(self, plot_realtime=0, save=0):
        plot_realtime = 0
        self.clearPlot()
        n = self.simulationsMenegement.active_experiments
        if not n and plot_realtime == 0:
            self.replot_experimental()
            return 0
        elif plot_realtime == 1:
            return
        else:
            for n in self.simulationsMenegement.active_experiments:
                xp = self.simulationsMenegement.simulations_list["results"][n]
                self.widget_layer_params.table_layer_params.setData(
                    self.simulationsMenegement.simulations_list["layer_params"]
                    [n].to_numpy())
                self.widget_layer_params.table_layer_params.setHorizontalHeaderLabels(
                    self.simulationsMenegement.simulations_list["layer_params"]
                    [n].columns)
                self.widget_layer_params.table_stimulus_params.setData(
                    self.simulationsMenegement.
                    simulations_list["simulation_params"][n].to_numpy())
                self.widget_layer_params.table_stimulus_params.setHorizontalHeaderLabels(
                    self.simulationsMenegement.
                    simulations_list["simulation_params"][n].columns)
                self.replot_all(xp, save=save)
        self.replot_experimental()

    def update(self):
        for q in self.ports:
            try:
                prog, self.exp = q.get(block=False)
                self.ctrLayout.progress.setValue(prog)
                self.replot_results(plot_realtime=1)
            except Queue.Empty:
                pass

    def simple_update(self):
        try:
            self.replot_results(plot_realtime=1)
        except Queue.Empty:
            pass


class LayerStructure():
    def __init__(self, sim_num):
        for val in [
                "Ms", "Ku", "J", "th", "alpha", "AMR", "SMR", "Rx0", "Ry0",
                "w", "l"
        ]:
            setattr(
                self, val,
                np.asarray(
                    plotter.simulationsMenegement.
                    simulations_list["layer_params"][sim_num][val].values,
                    dtype=np.float32))
        self.Kdir = self.get_kdir(
            plotter.simulationsMenegement.simulations_list["layer_params"]
            [sim_num]["Kdir"].values)
        self.Ndemag2 = self.get_Ndemag(
            plotter.simulationsMenegement.simulations_list["layer_params"]
            [sim_num]["N"].values)
        self.number_of_layers = len(self.Ms)

    def get_kdir(self, value):
        listOfParams = []
        for n in range(0, len(value)):
            tmp = value[n]
            tmp = tmp.replace("[", "").replace("]", "")
            v_tmp = (tmp.split(" "))
            res = np.array(list(map(float, v_tmp)))
            res = normalize(res)
            listOfParams.append(res)
        return listOfParams

    def get_Ndemag(self, value):
        tmp = plotter.widget_layer_params.table_layer_params.item(0, 7).text()
        tmp = tmp.replace("[", "").replace("]", "")
        v_tmp = (tmp.split(" "))
        res = np.array(list(map(float, v_tmp)))
        N = np.array([[res[0], 0, 0], [0, res[1], 0], [0, 0, res[2]]])
        return N


class SimulationStimulus():
    def __init__(self, sim_num):
        data = plotter.simulationsMenegement.simulations_list[
            "simulation_params"][sim_num]

        # self.theta = np.array(data["HTheta"].values[0], dtype=np.float32)
        self.back = np.array(data["Hback"].values[0], dtype=np.int)
        # self.phi = np.array(data["HPhi"].values[0], dtype=np.float32)

        if data["H"].values[1] != "-" and data["HPhi"].values[
                1] == "-" and data["HTheta"].values[1] == "-":
            self.mode = "H"
            self.STEPS = np.array(data["H"].values[1], dtype=np.float32)
            self.Hmin = np.array(data["H"].values[0], dtype=np.float32)
            self.Hmax = np.array(data["H"].values[2], dtype=np.float32)
            self.ThetaMin = np.array(data["HTheta"].values[0],
                                     dtype=np.float32)
            self.ThetaMax = np.array(data["HTheta"].values[0],
                                     dtype=np.float32)
            self.PhiMin = np.array(data["HPhi"].values[0], dtype=np.float32)
            self.PhiMax = np.array(data["HPhi"].values[0], dtype=np.float32)

        elif data["HPhi"].values[1] != "-" and data["H"].values[
                1] == "-" and data["HTheta"].values[1] == "-":
            self.mode = "phi"
            self.STEPS = np.array(data["HPhi"].values[1], dtype=np.float32)
            self.Hmin = np.array(data["H"].values[0], dtype=np.float32)
            self.Hmax = np.array(data["H"].values[0], dtype=np.float32)
            self.ThetaMin = np.array(data["HTheta"].values[0],
                                     dtype=np.float32)
            self.ThetaMax = np.array(data["HTheta"].values[0],
                                     dtype=np.float32)
            self.PhiMin = np.array(data["HPhi"].values[0], dtype=np.float32)
            self.PhiMax = np.array(data["HPhi"].values[2], dtype=np.float32)
        elif data["HTheta"].values[1] != "-" and data["H"].values[
                1] == "-" and data["HPhi"].values[1] == "-":
            self.mode = "theta"
            self.STEPS = np.array(data["HTheta"].values[1], dtype=np.float32)
            self.Hmin = np.array(data["H"].values[0], dtype=np.float32)
            self.Hmax = np.array(data["H"].values[0], dtype=np.float32)
            self.ThetaMin = np.array(data["HTheta"].values[0],
                                     dtype=np.float32)
            self.ThetaMax = np.array(data["HTheta"].values[2],
                                     dtype=np.float32)
            self.PhiMin = np.array(data["HPhi"].values[0], dtype=np.float32)
            self.PhiMax = np.array(data["HPhi"].values[0], dtype=np.float32)
        else:
            print("Stimulus error")
        self.H_sweep, self.Hmag = get_stimulus2(self.Hmin, self.Hmax,
                                                self.ThetaMin, self.ThetaMax,
                                                self.PhiMin, self.PhiMax,
                                                self.STEPS, self.back,
                                                self.mode)
        self.fmin = np.array(data["f"].values[0], dtype=np.float32)
        self.fsteps = np.array(data["f"].values[1], dtype=np.int)
        self.fmax = np.array(data["f"].values[2], dtype=np.float32)
        self.LLGtime = np.array(data["LLGtime"].values[0], dtype=np.float32)
        self.LLGsteps = int(
            np.array(data["LLGsteps"].values[0], dtype=np.float32))
        self.freqs = np.linspace(self.fmin, self.fmax, self.fsteps)
        self.spectrum_len = (self.LLGsteps) // 2
        self.PIMM_delta_f = 1 / self.LLGtime
        self.fphase = np.array(data["fphase"].values[0], dtype=np.float32)


class SimulationResults():
    def __init__(self, Stimulus, SpinDevice):
        self.H = []
        self.Rx = []
        self.Ry = []
        self.Rz = []
        self.Hmag_out = []
        self.Mlayers = np.empty((0, SpinDevice.number_of_layers, 3), float)
        self.M_avg = np.empty((0, 3), float)
        self.R_net = np.empty((0, 3), float)
        self.Spectrogram_data = np.empty((0, Stimulus.spectrum_len), float)
        self.Spectrogram_VSD = np.empty((0, len(Stimulus.freqs)), float)


def Simulate():
    runner = SimulationRunner()
    while 1 == 1:
        global stop
        if stop == 0:
            list_todo = plotter.simulation_manager.active_experiments
            plotter.simulation_manager.active_experiments = []

            backend = plotter.ctrLayout.Backend_choose.currentText()
            print("Simulation started with backend", backend)

            for sim_num in list_todo:
                spin_device = LayerStructure(sim_num)
                stimulus = SimulationStimulus(sim_num)
                sim_name = plotter.simulation_manager.simulations_list[
                    "settings"][sim_num][1]
                SimulationTimeStamp = datetime.datetime.now()
                sim_results = SimulationResults(Stimulus=stimulus,
                                                SpinDevice=spin_device)

                runner.run_simulation(spin_device, stimulus, sim_results,
                                      plotter, sim_num, curve, sim_name)

        else:
            time.sleep(0.001)


if __name__ == "__main__":
    app = QtGui.QApplication([])
    MainWindow = QtGui.QMainWindow()
    QtCore.QCoreApplication.setOrganizationName("PyMag")
    QtCore.QCoreApplication.setApplicationName(PyMagVersion)
    plotter = Ui_MainWindow()
    curve = plotter.getPort()
    if (sys.flags.interactive != 1) or not hasattr(QtCore, 'PYQT_VERSION'):
        QtGui.QApplication.instance().exec_()
    sys.exit(app.exec_())
